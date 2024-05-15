import shutil
import os
import datetime
from datetime import timedelta


try:
    from pptx import Presentation
except:
    os.system("pip install python-pptx")
    try:
        from pptx import Presentation
    except:
        os.system("C:\\Users\\dev\\AppData\\Local\\Programs\\Python\\Python312\\python.exe -m pip install python-pptx")
        from pptx import Presentation
        

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.table import Table
import sys

def table_contents(table, thursday):
    # 슬라이드 내의 테이블 찾기
    # 테이블 내의 셀 내용 가져오기
    # for row in table.rows:
    #     for cell in row.cells:
    #         for i,paragraph in enumerate(cell.text_frame.paragraphs):
    #             print(paragraph.text)
    # table.cell(0, 1).text = text
    # 목요일 기준
    thismonday = thursday - timedelta(days=3)
    thisfriday = thursday + timedelta(days=1)
    formattedMonday1 = thismonday.strftime("%m/%d")
    formattedFriday1 = thisfriday.strftime("%m/%d")
    nextmonday = thursday + timedelta(days=4)
    nextfriday = thursday + timedelta(days=8)
    formattedMonday2 = nextmonday.strftime("%m/%d")
    formattedFriday2 = nextfriday.strftime("%m/%d")
    arr_tit = []
    arr_tit.append(f'금주 실적 ({formattedMonday1}~{formattedFriday1})')
    arr_tit.append(f'차주 계획 ({formattedMonday2}~{formattedFriday2})')
    for i in range(2):
        # text_frame1 = table.cell(0, 1).text_frame
        text_frame1 = table.cell(0, i+1).text_frame

        for paragraph in text_frame1.paragraphs:
             # 텍스트 정렬을 가운데로 설정
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = '맑은 고딕' # 폰트 이름 설정
                font.size = Pt(11)  # 폰트 크기 설정
                font.bold = True
                run.text = arr_tit[i]
    

def main():
    targetFilePath=''
    try:
        # 현재 날짜 및 시간을 가져옴
        currentDate = datetime.datetime.now()
        currentDay = datetime.datetime.now().weekday()
        formatted_thursday = ''
        if currentDay < 3 :
            print("이번주 보고서 제출날인 목요일 전입니다.")
            until_thursday = (3-currentDay)%7
            this_thursday = currentDate + datetime.timedelta(days=until_thursday)
            formatted_thursday = this_thursday.strftime("%Y. %m. %d")
        elif currentDay == 3:
            print("목요일 입니다")
            formatted_thursday = currentDate.strftime("%Y. %m. %d")
            this_thursday = currentDate
        elif currentDay > 3:
            print("목요일 이후입니다.")
            until_thursday = (3-currentDay)%7
            this_thursday = currentDate + datetime.timedelta(days=until_thursday)
            formatted_thursday = this_thursday.strftime("%Y. %m. %d")
        # 복사할 PowerPoint 파일의 경로와 이름
        sourceFilePath = r'C:\Users\dev\Desktop\양식\주간보고_.pptx'

    #     # 파일을 복사할 대상 폴더 경로
        targetFolderPath = r"""C:\Users\dev\Desktop\주간보고서"""

    #     # 새로운 파일 이름 설정 (원하는 파일명으로 변경)
        currentDate = datetime.datetime.now()
        fileName = f"주간보고서_이기원_{this_thursday.strftime('%Y%m%d')}.pptx"

    #     # 대상 폴더에 저장할 새로운 파일의 전체 경로
        targetFilePath = os.path.join(targetFolderPath, fileName)

        # PowerPoint 파일을 복사하여 새로운 이름으로 저장
        shutil.copy(sourceFilePath, targetFilePath)

        print("보고서 생성 완료")
    except Exception as ex:
        print("에러 발생:", ex)

    try:
        
        pptFilePath = targetFilePath
    
        # pptFilePath = """C:\Users\dev\Desktop\주간보고서\주간보고서_이기원_20240425"""
        

        #파워포인트 변수 생성
        prs = Presentation(pptFilePath)

        #첫번째 슬라이드 (마스터 폼)
        silde = prs.slides[0]

        shape_list = silde.shapes
        shape_idx = {}

        #{shape이름:순번}의 형태로 dict 생성
        #{ppt에서 설정한 이름1 : 0,ppt에서 설정한 이름2 : 1}
        for idx,value in enumerate(shape_list):
            shape_idx[value.name]=idx

        #첫번째 슬라이드 날짜 : titleDate
        #titleDate란 텍스트 박스 shape 가져오기
        shape = shape_list[shape_idx['titleDate']]

        tf = shape.text_frame
        #텍스트 박스 내 글자 삭제
        print(tf.text)
        tf.clear()
        #택스트 박스 문단 선택
        para = tf.paragraphs[0]
        #중앙 정렬
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()



        # # 날짜 포맷팅 (YYYY. MM. DD)
        # formattedDate = currentDate.strftime("%Y. %m. %d")
        #택스트 내용
        run.text = formatted_thursday
        print('날짜 변경 완')

        #두번째 슬라이드 선택
        silde = prs.slides[1]
        shape_list = silde.shapes
        shape_idx = {}

        #{shape이름:순번}의 형태로 dict 생성
        #{ppt에서 설정한 이름1 : 0,ppt에서 설정한 이름2 : 1}
        for idx,value in enumerate(shape_list):
            shape_idx[value.name]=idx

        #첫번째 슬라이드 날짜 : titleDate
        #titleDate란 텍스트 박스 shape 가져오기

        shape = shape_list[shape_idx['report_table']]

        tbl=shape.table
        table_contents(tbl, this_thursday)
        #파일 저장
        prs.save(pptFilePath)

    except Exception as ex:
        print("에러 발생:", ex)

    # input("Press Enter to close the window...")
    



if __name__ == "__main__":
    main()

    # user_input = input("아무 키나 입력하세요 (종료됩니다): ")
    # sys.exit()


